from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
PPTX_PATH = BASE_DIR / "Pothole_Detection_Committee_Presentation.pptx"
DB_IMAGE = BASE_DIR.parent / "report-assets" / "png" / "figure-5-er-schema.png"
CONFUSION_IMAGE = BASE_DIR.parent / "ai-model" / "trained_model" / "confusion_matrix.png"
RESULTS_IMAGE = BASE_DIR.parent / "ai-model" / "trained_model" / "results.png"

PRIMARY = RGBColor(12, 51, 77)
SECONDARY = RGBColor(29, 111, 141)
BG = RGBColor(246, 248, 251)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(92, 107, 122)
WHITE = RGBColor(255, 255, 255)


def set_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.45)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.68), Inches(1.22), Inches(12), Inches(0.5))
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Aptos"
        sr.font.size = Pt(13)
        sr.font.color.rgb = MUTED


def add_bullets(slide, items, left, top, width, height, font_size=18, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(9)
        p.bullet = True


def add_label_box(slide, title, lines, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = SECONDARY
    shape.line.width = Pt(1.1)

    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.3), Inches(0.35))
    ttf = title_box.text_frame
    tp = ttf.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.name = "Aptos"
    tr.font.size = Pt(14)
    tr.font.bold = True
    tr.font.color.rgb = SECONDARY

    body_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.45), width - Inches(0.35), height - Inches(0.55))
    btf = body_box.text_frame
    btf.word_wrap = True
    for idx, line in enumerate(lines):
        p = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        p.space_after = Pt(5)


def add_image_contain(slide, image_path, left, top, width, height):
    with Image.open(image_path) as img:
        img_w, img_h = img.size

    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        final_w = width
        final_h = width / img_ratio
        final_left = left
        final_top = top + (height - final_h) / 2
    else:
        final_h = height
        final_w = height * img_ratio
        final_top = top
        final_left = left + (width - final_w) / 2

    slide.shapes.add_picture(str(image_path), final_left, final_top, width=final_w, height=final_h)


def add_database_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide)
    add_title(slide, "Database Design", "Core entities and spatial data support behind the reporting workflow")
    add_image_contain(slide, DB_IMAGE, Inches(0.7), Inches(1.7), Inches(7.2), Inches(4.9))
    add_bullets(
        slide,
        [
            "PostgreSQL stores users, vehicles, pothole reports, and related operational data.",
            "PostGIS adds geospatial support so pothole locations can be mapped and queried efficiently.",
            "The schema is designed to support report tracking, status updates, and dashboard filtering.",
            "This database layer connects the mobile reporting flow with the administrative dashboard.",
        ],
        Inches(8.15), Inches(1.95), Inches(4.2), Inches(3.4), 17,
    )
    add_label_box(
        slide,
        "Why It Matters",
        [
            "Reliable persistence for reports and metadata",
            "Spatial queries for map-based monitoring",
            "Clear structure for future scaling and analytics",
        ],
        Inches(8.1), Inches(5.15), Inches(4.25), Inches(1.3),
    )


def add_ai_charts_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_top_band(slide)
    add_title(slide, "AI Results Charts", "Visual evidence from model evaluation and training history")

    add_label_box(slide, "Confusion Matrix", ["Shows how well the model separates pothole detections from background and errors."], Inches(0.8), Inches(1.75), Inches(5.8), Inches(0.95))
    add_image_contain(slide, CONFUSION_IMAGE, Inches(0.8), Inches(2.75), Inches(5.8), Inches(3.6))

    add_label_box(slide, "Training Curves", ["Shows loss reduction and metric improvement during training, supporting model convergence."], Inches(6.75), Inches(1.75), Inches(5.8), Inches(0.95))
    add_image_contain(slide, RESULTS_IMAGE, Inches(6.75), Inches(2.75), Inches(5.8), Inches(3.6))

    footer = slide.shapes.add_textbox(Inches(0.9), Inches(6.45), Inches(11.6), Inches(0.4))
    ftf = footer.text_frame
    p = ftf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "These charts complement the summary metrics by showing class behavior and the model's training stability."
    r.font.name = "Aptos"
    r.font.size = Pt(14)
    r.font.color.rgb = MUTED


def main():
    prs = Presentation(PPTX_PATH)
    add_database_slide(prs)
    add_ai_charts_slide(prs)
    prs.save(PPTX_PATH)
    print(f"Updated presentation: {PPTX_PATH}")


if __name__ == "__main__":
    main()
