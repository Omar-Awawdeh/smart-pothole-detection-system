from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
ASSETS_DIR = BASE_DIR.parent / "report-assets" / "png"
OUTPUT_PATH = BASE_DIR / "Pothole_Detection_Committee_Presentation.pptx"

PRIMARY = RGBColor(12, 51, 77)
SECONDARY = RGBColor(29, 111, 141)
ACCENT = RGBColor(231, 111, 81)
GOLD = RGBColor(245, 180, 51)
BG = RGBColor(246, 248, 251)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(92, 107, 122)
WHITE = RGBColor(255, 255, 255)


def set_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.45)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()


def add_title(slide, title, subtitle=None, dark=False):
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE if dark else PRIMARY
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.68), Inches(1.22), Inches(12), Inches(0.5))
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Aptos"
        sr.font.size = Pt(13)
        sr.font.color.rgb = WHITE if dark else MUTED


def add_bullets(slide, items, left, top, width, height, font_size=20, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(9)
        p.bullet = True
    return box


def add_label_value_box(slide, title, body, left, top, width, height, fill_color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = SECONDARY
    shape.line.width = Pt(1.2)

    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.3), Inches(0.3))
    title_tf = title_box.text_frame
    tp = title_tf.paragraphs[0]
    tr = tp.add_run()
    tr.text = title
    tr.font.name = "Aptos"
    tr.font.bold = True
    tr.font.size = Pt(14)
    tr.font.color.rgb = SECONDARY

    body_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.45), width - Inches(0.35), height - Inches(0.55))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    for idx, line in enumerate(body):
        p = body_tf.paragraphs[0] if idx == 0 else body_tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        p.space_after = Pt(5)
    return shape


def add_image_contain(slide, image_path, left, top, width, height):
    from PIL import Image

    with Image.open(image_path) as img:
        img_w, img_h = img.size
    box_ratio = width / height
    img_ratio = img_w / img_h

    if img_ratio > box_ratio:
        final_w = width
        final_h = width / img_ratio
        final_left = left
        final_top = top + (height - final_h) / 2
    else:
        final_h = height
        final_w = height * img_ratio
        final_top = top
        final_left = left + (width - final_w) / 2

    slide.shapes.add_picture(str(image_path), final_left, final_top, width=final_w, height=final_h)


def metric_card(slide, value, label, left, top, fill_color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(2.45), Inches(1.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    value_box = slide.shapes.add_textbox(left, top + Inches(0.22), Inches(2.45), Inches(0.5))
    vtf = value_box.text_frame
    vp = vtf.paragraphs[0]
    vp.alignment = PP_ALIGN.CENTER
    vr = vp.add_run()
    vr.text = value
    vr.font.name = "Aptos Display"
    vr.font.size = Pt(22)
    vr.font.bold = True
    vr.font.color.rgb = WHITE

    label_box = slide.shapes.add_textbox(left, top + Inches(0.82), Inches(2.45), Inches(0.35))
    ltf = label_box.text_frame
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    lr.text = label
    lr.font.name = "Aptos"
    lr.font.size = Pt(12)
    lr.font.color.rgb = WHITE


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_background(slide, PRIMARY)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(6.75), Inches(13.333), Inches(0.75))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    add_title(slide, "Smart Pothole Detection System", "Graduation Project Committee Presentation", dark=True)
    team = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(6.5), Inches(2.2))
    tf = team.text_frame
    tf.word_wrap = True
    lines = [
        "Computer Engineering - Yarmouk University",
        "Team Members: Omar Awawdeh, Hamza Al-Safi",
        "Academic Year: 2025-2026",
        "Committee Presentation",
    ]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(20 if idx == 0 else 18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(7)
    add_image_contain(slide, ASSETS_DIR / "figure-1-high-level.png", Inches(7.2), Inches(1.45), Inches(5.35), Inches(4.9))

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_band(slide)
    add_title(slide, "Problem Statement", "Why a smart pothole monitoring workflow is needed")
    add_bullets(
        slide,
        [
            "Manual road inspection is slow, costly, and difficult to keep up-to-date across large areas.",
            "Potholes create safety risks, vehicle damage, and higher maintenance costs when they are reported late.",
            "Citizens may report issues, but reports are often inconsistent, duplicated, or missing precise location data.",
            "Road authorities need timely, GPS-tagged evidence to prioritize maintenance more effectively.",
        ],
        Inches(0.8), Inches(1.7), Inches(6.1), Inches(4.8), 20,
    )
    add_label_value_box(
        slide,
        "Core Idea",
        [
            "Use a smartphone camera and on-device AI to detect potholes in real time, attach GPS coordinates,",
            "and send structured reports to a central dashboard for verification and tracking.",
        ],
        Inches(7.2), Inches(2.0), Inches(5.2), Inches(2.0),
    )
    add_label_value_box(
        slide,
        "Why On Device?",
        [
            "Lower latency",
            "Works with limited connectivity",
            "Lower bandwidth and cloud cost",
            "Better privacy than streaming video to the cloud",
        ],
        Inches(7.2), Inches(4.25), Inches(5.2), Inches(1.9),
    )

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_band(slide)
    add_title(slide, "Objectives", "Target outcomes for the complete system")
    goals = [
        ("Detect", "Identify potholes from live camera frames in real time on Android."),
        ("Locate", "Capture GPS coordinates for each valid detection."),
        ("Report", "Upload structured reports with images and metadata to the backend."),
        ("Visualize", "Provide a dashboard with map, records, and status management."),
        ("Deploy", "Run the full system on a VPS with a production-style architecture."),
    ]
    x = Inches(0.8)
    for title, body in goals:
        add_label_value_box(slide, title, [body], x, Inches(1.8), Inches(2.35), Inches(3.2), fill_color=WHITE)
        x += Inches(2.5)

    note = slide.shapes.add_textbox(Inches(1.0), Inches(5.55), Inches(11.3), Inches(0.9))
    ntf = note.text_frame
    p = ntf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "This project is designed as an end-to-end workflow, not only an isolated AI model."
    r.font.name = "Aptos"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = PRIMARY

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_band(slide)
    add_title(slide, "Proposed Solution", "End-to-end workflow from detection to maintenance visibility")
    add_image_contain(slide, ASSETS_DIR / "figure-1-high-level.png", Inches(0.7), Inches(1.55), Inches(7.35), Inches(5.3))
    add_bullets(
        slide,
        [
            "Android app captures live frames and runs YOLOv8n through TensorFlow Lite.",
            "Valid detections are filtered, geotagged, deduplicated locally, and queued for upload.",
            "Backend stores pothole records in PostgreSQL/PostGIS and manages report lifecycle.",
            "Dashboard allows authorities to review locations, images, and repair status.",
        ],
        Inches(8.25), Inches(1.95), Inches(4.3), Inches(4.4), 18,
    )

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_band(slide)
    add_title(slide, "System Architecture", "Current implementation stack in the repository")
    add_label_value_box(slide, "Mobile Layer", ["Kotlin", "Jetpack Compose", "CameraX", "TensorFlow Lite", "GPS + WorkManager"], Inches(0.8), Inches(1.7), Inches(2.35), Inches(3.7))
    add_label_value_box(slide, "AI Layer", ["YOLOv8n object detector", "Float16 TFLite export", "Recall-tuned thresholds"], Inches(3.35), Inches(1.7), Inches(2.35), Inches(3.7))
    add_label_value_box(slide, "Backend Layer", ["ASP.NET Core 8 Web API", "Entity Framework Core", "REST endpoints"], Inches(5.9), Inches(1.7), Inches(2.35), Inches(3.7))
    add_label_value_box(slide, "Data Layer", ["PostgreSQL 16", "PostGIS spatial support", "Cloudflare R2 image storage"], Inches(8.45), Inches(1.7), Inches(2.35), Inches(3.7))
    add_label_value_box(slide, "Dashboard + Ops", ["React + Vite + Tailwind", "React Leaflet map", "Docker Compose + NGINX on VPS"], Inches(11.0), Inches(1.7), Inches(1.55), Inches(3.7))
    footer = slide.shapes.add_textbox(Inches(0.9), Inches(5.9), Inches(11.7), Inches(0.6))
    ftf = footer.text_frame
    p = ftf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Current runtime stack documented in docs/07-current-implementation.md and docs/05-deployment.md"
    r.font.name = "Aptos"
    r.font.size = Pt(15)
    r.font.color.rgb = MUTED

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_band(slide)
    add_title(slide, "AI Model and Dataset", "Why YOLOv8n is suitable for mobile pothole detection")
    add_bullets(
        slide,
        [
            "Model: YOLOv8n (nano), selected for a strong speed/accuracy trade-off on mobile devices.",
            "Dataset: 2,642 images and 9,077 labeled potholes collected from public datasets.",
            "Data split: 80.6% train, 14.5% validation, 5.0% test.",
            "Deployment format: float16 TFLite model, about 6 MB, suitable for Android packaging.",
        ],
        Inches(0.85), Inches(1.7), Inches(6.2), Inches(4.5), 18,
    )
    metric_card(slide, "2,642", "Images", Inches(7.35), Inches(1.8), SECONDARY)
    metric_card(slide, "9,077", "Annotations", Inches(10.0), Inches(1.8), PRIMARY)
    metric_card(slide, "~6 MB", "TFLite Size", Inches(7.35), Inches(3.55), ACCENT)
    metric_card(slide, "25-30 FPS", "Expected Mobile Speed", Inches(10.0), Inches(3.55), RGBColor(54, 122, 79))
    small = slide.shapes.add_textbox(Inches(7.3), Inches(5.45), Inches(4.95), Inches(0.7))
    stf = small.text_frame
    sp = stf.paragraphs[0]
    sp.text = "Training used recall tuning to improve the chance of catching more potholes in real-world use."
    sp.font.name = "Aptos"
    sp.font.size = Pt(15)
    sp.font.color.rgb = TEXT

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_band(slide)
    add_title(slide, "Android Application", "Main features implemented in the mobile app")
    add_bullets(
        slide,
        [
            "Live camera preview with real-time pothole detection overlay.",
            "On-device inference using TensorFlow Lite for low-latency performance.",
            "GPS integration to tag each valid detection with location metadata.",
            "History and settings screens for review and configuration.",
            "Background upload queue for reliable synchronization when connectivity changes.",
        ],
        Inches(0.8), Inches(1.75), Inches(6.0), Inches(4.6), 18,
    )
    add_label_value_box(slide, "Engineering Focus", ["Low latency", "Offline-friendly flow", "Mobile-first UX", "Reliable upload behavior"], Inches(7.2), Inches(1.95), Inches(5.1), Inches(2.3), fill_color=WHITE)
    add_label_value_box(slide, "Representative Screens", ["Detection", "History", "Settings", "Debug / verification support"], Inches(7.2), Inches(4.5), Inches(5.1), Inches(1.75), fill_color=WHITE)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_band(slide)
    add_title(slide, "Backend and Dashboard", "How reports are stored, reviewed, and managed")
    add_image_contain(slide, ASSETS_DIR / "figure-3-dashboard.png", Inches(0.7), Inches(1.6), Inches(7.2), Inches(5.2))
    add_bullets(
        slide,
        [
            "Backend uses ASP.NET Core 8 with PostgreSQL/PostGIS for API and spatial data handling.",
            "Dashboard uses React, Vite, Tailwind, and React Leaflet for map-based monitoring.",
            "Authorities can browse pothole records, inspect images, and update status such as verified or repaired.",
            "Map view loads all pothole points and auto-fits bounds for easier monitoring.",
        ],
        Inches(8.15), Inches(1.85), Inches(4.45), Inches(4.7), 17,
    )

    # Slide 9
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_band(slide)
    add_title(slide, "Results", "Best validation metrics and deployment operating point")
    metric_card(slide, "81.45%", "mAP@50", Inches(0.8), Inches(1.75), PRIMARY)
    metric_card(slide, "52.13%", "mAP@50-95", Inches(3.5), Inches(1.75), SECONDARY)
    metric_card(slide, "79.98%", "Precision", Inches(6.2), Inches(1.75), ACCENT)
    metric_card(slide, "75.90%", "Recall", Inches(8.9), Inches(1.75), RGBColor(54, 122, 79))
    add_label_value_box(
        slide,
        "Threshold Sweep Recommendation",
        [
            "Confidence threshold: 0.30",
            "NMS IoU threshold: 0.45",
            "Precision / Recall at recommended operating point: 77.41% / 77.02%",
            "Recall improved by +3.86 points over the earlier baseline.",
        ],
        Inches(0.9), Inches(3.65), Inches(5.85), Inches(2.15), fill_color=WHITE,
    )
    add_label_value_box(
        slide,
        "Interpretation",
        [
            "The tuned model keeps precision near 80% while reducing missed potholes.",
            "This trade-off is suitable for a field reporting system where missing defects is costly.",
        ],
        Inches(7.0), Inches(3.65), Inches(5.35), Inches(2.15), fill_color=WHITE,
    )

    # Slide 10
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_band(slide)
    add_title(slide, "Deployment and Demo Readiness", "Production-style setup used for the integrated system")
    add_image_contain(slide, ASSETS_DIR / "figure-4-deployment.png", Inches(0.7), Inches(1.6), Inches(6.75), Inches(5.2))
    add_bullets(
        slide,
        [
            "System is deployed on a Linux VPS using Docker Compose and NGINX.",
            "Dashboard domain: potholesystem.tech; API domain: api.potholesystem.tech.",
            "PostgreSQL + PostGIS provides persistent geospatial storage.",
            "Uploaded images can be served from Cloudflare R2 through a public image domain.",
        ],
        Inches(7.8), Inches(1.85), Inches(4.45), Inches(4.8), 17,
    )

    # Slide 11
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_top_band(slide)
    add_title(slide, "Limitations and Future Work", "What can be improved after the current prototype")
    add_label_value_box(slide, "Current Limitations", ["Recall is improved but some potholes are still missed.", "Training data is mostly public and may not fully match local Jordan roads.", "Performance can vary with lighting, glare, weather, and motion blur.", "Current model detects potholes only, without severity grading."], Inches(0.8), Inches(1.8), Inches(5.9), Inches(3.8), fill_color=WHITE)
    add_label_value_box(slide, "Future Work", ["Collect local road data and fine-tune on regional conditions.", "Add severity classification and other road-hazard classes.", "Validate the tuned TFLite export more extensively on target devices.", "Expand robustness for night-time and adverse weather conditions."], Inches(6.95), Inches(1.8), Inches(5.45), Inches(3.8), fill_color=WHITE)

    # Slide 12
    slide = prs.slides.add_slide(blank)
    set_background(slide, PRIMARY)
    add_title(slide, "Conclusion", "An end-to-end AI-assisted road monitoring solution", dark=True)
    add_bullets(
        slide,
        [
            "The project delivers a complete workflow: detection, geotagging, reporting, storage, and dashboard monitoring.",
            "The mobile AI model achieves strong academic results while remaining deployable on Android devices.",
            "The integrated architecture demonstrates practical value for municipalities and road maintenance teams.",
            "This work can be extended into a larger smart-city road condition monitoring platform.",
        ],
        Inches(0.9), Inches(1.75), Inches(7.0), Inches(4.2), 20, WHITE,
    )
    thanks = slide.shapes.add_textbox(Inches(8.2), Inches(2.35), Inches(4.0), Inches(1.2))
    ttf = thanks.text_frame
    p = ttf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Thank You"
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = GOLD
    p2 = ttf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Questions and Discussion"
    r2.font.name = "Aptos"
    r2.font.size = Pt(18)
    r2.font.color.rgb = WHITE

    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_deck()
    print(f"Created presentation: {OUTPUT_PATH}")
